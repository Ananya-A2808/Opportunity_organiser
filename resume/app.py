from flask import Flask, render_template, request, send_file, url_for, jsonify
import os
import pdfkit
from werkzeug.utils import secure_filename
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = Flask(__name__)

# Configure upload folder
UPLOAD_FOLDER = 'static/uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Create upload folder if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Path to wkhtmltopdf executable
config = pdfkit.configuration(
    wkhtmltopdf=r"/usr/local/bin/wkhtmltopdf"
)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def create_presentation(data):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data['name']
    subtitle.text = data['title']
    
    # About slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "About Me"
    tf = body_shape.text_frame
    tf.text = data['about']
    
    # Education slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Education"
    tf = body_shape.text_frame
    
    if data.get('edu1'):
        p = tf.add_paragraph()
        p.text = data['edu1']
    if data.get('edu2'):
        p = tf.add_paragraph()
        p.text = data['edu2']
    if data.get('edu3'):
        p = tf.add_paragraph()
        p.text = data['edu3']
    
    # Experience slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Experience"
    tf = body_shape.text_frame
    
    if data.get('exp1'):
        p = tf.add_paragraph()
        p.text = data['exp1']
    if data.get('exp2'):
        p = tf.add_paragraph()
        p.text = data['exp2']
    
    # Skills slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Skills"
    tf = body_shape.text_frame
    
    if data.get('skills'):
        for skill in data['skills'].split(','):
            p = tf.add_paragraph()
            p.text = skill.strip()
    
    # References slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "References"
    tf = body_shape.text_frame
    
    if data.get('ref1'):
        p = tf.add_paragraph()
        p.text = data['ref1']
    if data.get('ref2'):
        p = tf.add_paragraph()
        p.text = data['ref2']
    
    # Save the presentation
    presentation_path = "resume_presentation.pptx"
    prs.save(presentation_path)
    return presentation_path

@app.route("/", methods=["GET", "POST"])
def form():
    if request.method == "POST":
        data = request.form.to_dict()
        
        # Handle image upload
        if 'profile_image' in request.files:
            file = request.files['profile_image']
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                data['profile_image'] = url_for('static', filename=f'uploads/{filename}')
        
        # Choose template based on format
        template = "resume_professional.html" if data.get('format') == 'professional' else "resume.html"
        
        return render_template(template, data=data)
    return render_template("form.html")

@app.route("/download", methods=["POST"])
def download():
    data = request.form.to_dict()
    
    # Handle image in the data
    if 'profile_image' in data:
        # Convert image to base64 for PDF embedding
        image_path = os.path.join(app.root_path, data['profile_image'].lstrip('/'))
        if os.path.exists(image_path):
            with open(image_path, 'rb') as img_file:
                data['profile_image'] = f"data:image/jpeg;base64,{base64.b64encode(img_file.read()).decode()}"
    
    # Choose template based on format
    template = "resume_professional.html" if data.get('format') == 'professional' else "resume.html"
    html = render_template(template, data=data)

    # Enable local file access and set PDF options
    options = {
        'enable-local-file-access': '',
        'page-size': 'A4',
        'margin-top': '0.75in',
        'margin-right': '0.75in',
        'margin-bottom': '0.75in',
        'margin-left': '0.75in',
        'encoding': 'UTF-8',
        'no-outline': None,
        'quiet': ''
    }

    # Generate the PDF
    pdfkit.from_string(html, "resume.pdf", configuration=config, options=options)
    
    # If presentation format is requested
    if data.get('output_presentation'):
        presentation_path = create_presentation(data)
        return jsonify({
            'pdf_url': url_for('download_pdf'),
            'presentation_url': url_for('download_presentation')
        })
    
    return send_file("resume.pdf", as_attachment=True)

@app.route("/download_pdf")
def download_pdf():
    return send_file("resume.pdf", as_attachment=True)

@app.route("/download_presentation")
def download_presentation():
    return send_file("resume_presentation.pptx", as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
